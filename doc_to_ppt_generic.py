import streamlit as st
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF
import docx  # python-docx

st.set_page_config(page_title="Doc → PPT Converter", page_icon="🎯", layout="centered")

# --- Custom CSS ---
st.markdown("""
    <style>
        .main { background-color: green; }
        .stApp { background-color: #0f0f0f; color: white; }
        h1 { color: #ffffff; font-size: 2rem; }
        .stRadio label { color: #cccccc; }
        .stTextInput input { background-color: #1e1e1e; color: white; }
        .stButton button { background-color: #6c63ff; color: white; border-radius: 8px; }
    </style>
""", unsafe_allow_html=True)

# ─── DOCX → plain text ───────────────────────────────────────────────────────
def extract_docx_text(input_file):
    doc = docx.Document(input_file)
    return "\n".join([p.text for p in doc.paragraphs])

# ─── PDF → plain text ────────────────────────────────────────────────────────
def extract_pdf_text(input_file):
    doc = fitz.open(input_file)
    return "\n".join([page.get_text("text") for page in doc])

# ─── Extract images from DOCX ────────────────────────────────────────────────
def extract_docx_images(input_file, image_dir="docx_images"):
    doc = docx.Document(input_file)
    Path(image_dir).mkdir(exist_ok=True)
    images = []
    for i, rel in enumerate(doc.part._rels.values()):
        if "image" in rel.target_ref:
            img_data = rel.target_part.blob
            img_path = Path(image_dir) / f"docx_img_{i}.png"
            with open(img_path, "wb") as f:
                f.write(img_data)
            images.append(str(img_path))
    return images

# ─── Extract images from PDF ─────────────────────────────────────────────────
def extract_pdf_images(input_file, image_dir="pdf_images"):
    doc = fitz.open(input_file)
    Path(image_dir).mkdir(exist_ok=True)
    images = []
    for page_num, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            img_path = Path(image_dir) / f"pdf_img_{page_num}_{img_index}.png"
            if pix.n < 5:
                pix.save(img_path)
            else:
                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                pix1.save(img_path)
            pix = None
            images.append(str(img_path))
    return images

# ─── Split text into slides ───────────────────────────────────────────────────
def split_text(text, method, keyword=None):
    """
    method options:
      'heading'   – split on lines that look like headings (short, no period at end)
      'page'      – split on form-feed character \f (PDF page breaks)
      'keyword'   – split whenever a user-defined keyword appears at line start
      'paragraph' – split every N non-empty paragraphs (default 3)
    """
    slides = []

    if method == "heading":
        current = []
        for line in text.splitlines():
            stripped = line.strip()
            is_heading = (
                stripped
                and len(stripped) < 80
                and not stripped.endswith(".")
                and stripped[0].isupper()
            )
            if is_heading and current:
                slides.append("\n".join(current))
                current = [stripped]
            else:
                current.append(line)
        if current:
            slides.append("\n".join(current))

    elif method == "page":
        pages = text.split("\f")
        slides = [p.strip() for p in pages if p.strip()]

    elif method == "keyword" and keyword:
        current = []
        for line in text.splitlines():
            if line.strip().lower().startswith(keyword.strip().lower()) and current:
                slides.append("\n".join(current))
                current = [line]
            else:
                current.append(line)
        if current:
            slides.append("\n".join(current))

    elif method == "paragraph":
        paras = [l for l in text.splitlines() if l.strip()]
        chunk_size = 3
        for i in range(0, len(paras), chunk_size):
            slides.append("\n".join(paras[i:i+chunk_size]))

    return slides if slides else [text]

# ─── Build markdown from slides ──────────────────────────────────────────────
def build_markdown(slides):
    md_parts = []
    for i, slide in enumerate(slides):
        lines = slide.strip().splitlines()
        title = lines[0].strip() if lines else f"Slide {i+1}"
        body  = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
        md_parts.append(f"# {title}\n\n{body}")
    return "\n\n---\n\n".join(md_parts)

# ─── Markdown → PPTX via Pandoc ──────────────────────────────────────────────
def convert_md_to_pptx(md_file, output_file):
    cmd = ["pandoc", str(md_file), "-t", "pptx", "-o", str(output_file)]
    subprocess.run(cmd, check=True)

# ─── Style the PPTX ──────────────────────────────────────────────────────────
def style_pptx(pptx_file, images=None):
    prs = Presentation(pptx_file)
    for i, slide in enumerate(prs.slides):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 25)

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                shape.left  = Inches(0.4)
                shape.top   = Inches(0.4)
                shape.width = prs.slide_width - Inches(0.8)

        if images and i < len(images):
            try:
                slide.shapes.add_picture(
                    images[i], Inches(1), Inches(2.5),
                    width=Inches(4), height=Inches(2.8)
                )
            except Exception:
                pass

    prs.save(pptx_file)

# ─── Streamlit UI ─────────────────────────────────────────────────────────────
def main():
    st.title("📄 → 🎯 Doc to PPT Converter")
    st.caption("Upload any DOCX or PDF and choose how you want it split into slides.")

    uploaded_file = st.file_uploader("Upload your file", type=["docx", "pdf"])

    st.markdown("#### How do you want to split into slides?")
    method = st.radio(
        "Split method",
        options=["By Headings", "By Page Breaks", "By Custom Keyword", "By Paragraphs (every 3)"],
        label_visibility="collapsed"
    )

    keyword = None
    if method == "By Custom Keyword":
        keyword = st.text_input(
            "Enter your keyword (e.g. Q, Section, Chapter)",
            placeholder="e.g. Q or Chapter or Section"
        )
        st.caption("Each slide will start when this keyword appears at the beginning of a line.")

    output_name = st.text_input("Output filename", value="presentation.pptx")

    if st.button("Convert to PPTX 🚀") and uploaded_file:
        input_ext  = uploaded_file.name.split(".")[-1].lower()
        input_path = Path(f"input.{input_ext}")
        md_path    = Path("temp.md")
        output_path = Path(output_name if output_name.endswith(".pptx") else output_name + ".pptx")

        with open(input_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        method_map = {
            "By Headings":               "heading",
            "By Page Breaks":            "page",
            "By Custom Keyword":         "keyword",
            "By Paragraphs (every 3)":   "paragraph",
        }
        chosen = method_map[method]

        with st.spinner("Converting... please wait ⏳"):
            try:
                # Extract text
                if input_ext == "docx":
                    text   = extract_docx_text(input_path)
                    images = extract_docx_images(input_path)
                else:
                    text   = extract_pdf_text(input_path)
                    images = extract_pdf_images(input_path)

                # Split into slides
                slides = split_text(text, chosen, keyword)
                st.info(f"✅ {len(slides)} slide(s) detected using '{method}'")

                # Build markdown & convert
                md_content = build_markdown(slides)
                md_path.write_text(md_content, encoding="utf-8")
                convert_md_to_pptx(md_path, output_path)
                style_pptx(output_path, images)

                st.success("Done! Your presentation is ready 🎉")

                with open(output_path, "rb") as f:
                    st.download_button(
                        label="⬇️ Download PPTX",
                        data=f,
                        file_name=output_path.name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

            except Exception as e:
                st.error(f"Something went wrong: {e}")

if __name__ == "__main__":
    main()
